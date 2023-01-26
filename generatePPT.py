import configparser
import copy
import os
from datetime import datetime
from pptx import Presentation

config = configparser.ConfigParser()
config.read('config.ini')
templateDir = config['default']['templatePptDir']
outputDir = config['default']['resultPptDir']


def main():
    if not os.path.exists(templateDir):
        os.makedirs(templateDir)

    if not os.path.exists(outputDir):
        os.makedirs(outputDir)

    file_list = os.listdir(templateDir)

    txt_file = [file for file in file_list if file.endswith(".txt")]
    valid_pptx_file(file_list)
    valid_txt_file(txt_file)

    for file in txt_file:
        generate_ppt(file)


def valid_txt_file(txt_file):
    if not txt_file:
        err_msg = "{} {}".format("해당 경로에 txt 파일이 존재하지 않습니다.", os.path.abspath(templateDir))
        raise FileNotFoundError(err_msg)


def valid_pptx_file(file_list):
    pptx_file = [file for file in file_list if file.endswith(".pptx")]
    if not pptx_file:
        err_msg = "해당 경로에 pptx 파일이 존재하지 않습니다. {}".format(os.path.abspath(templateDir))
        raise FileNotFoundError(err_msg)
    if len(pptx_file) != 1:
        err_msg = "{} \n {} {}".format("pptx 파일은 한개만 존재해야 합니다.", os.path.abspath(templateDir), pptx_file)
        raise OSError(err_msg)


def generate_template_ppt_name():
    file_list = os.listdir(templateDir)
    for file in file_list:
        if file.endswith('.pptx') or file.endswith('.ppt'):
            return templateDir + '/' + file


def generate_output_file_name(file):
    file_name = file.strip('.txt')
    date_time = datetime.today().strftime('%Y%m%d-%H%M%S')
    return '{}/{}-{}.pptx'.format(outputDir, file_name, date_time)


def generate_ppt(file):
    try:
        template_ppt_name = generate_template_ppt_name()
        template_prs = Presentation(template_ppt_name)

        print("PPT 양식 파일 :", template_ppt_name)

        with open(templateDir + '/' + file, 'r') as f:
            reads = f.read()
            lyrics = reads.split('\n\n')  # TODO: 한 페이지를 구분하는 기준 -> config로 추출하기
            title = lyrics.pop(0)

            for entry in lyrics:
                duplicate_slide(template_prs, 0, title, entry)

        delete_slide(template_prs, 0)
        template_prs.save(generate_output_file_name(file))
        print("파일 생성 완료 :", generate_output_file_name(file))
    except Exception as e:
        print("파일 생성 실패. ", e)


def duplicate_slide(pres, index, title, entry):
    source = pres.slides[index]
    try:
        blank_slide_layout = pres.slide_layouts[6]
    except:
        blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts) - 1]

    new_slide = pres.slides.add_slide(blank_slide_layout)
    img_dict = {}

    for shp in source.shapes:
        if shp.shape_type == 13:
            # copy pictures
            with open(shp.name + '.jpg', 'wb') as f:
                f.write(shp.image.blob)

            img_dict[shp.name + '.jpg'] = [shp.left, shp.top, shp.width, shp.height]

            # add pictures
            for k, v in img_dict.items():
                pic = new_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
                pic.crop_right = shp.crop_right
                pic.crop_left = shp.crop_left
                pic.crop_top = shp.crop_top
                pic.crop_bottom = shp.crop_bottom
                os.remove(k)
        else:
            copy_shp = copy.deepcopy(shp)
            for paragraph in copy_shp.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text == '제목':
                        run.text = title
                    elif run.text == '가사':
                        run.text = entry

            newel = copy.deepcopy(copy_shp.element)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


if __name__ == "__main__":
    main()
